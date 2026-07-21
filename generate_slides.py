# /// script
# requires-python = ">=3.12,<3.13"
# dependencies = ["python-pptx==1.0.2"]
# ///
"""Generate the eight-slide TabFM Masterclass presentation from tutorial Markdown."""

from __future__ import annotations

import re
import sys
import tempfile
from collections.abc import Iterable, Sequence
from dataclasses import dataclass, field
from pathlib import Path
from typing import Final

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Inches, Pt

ROOT: Final = Path(__file__).resolve().parent
SOURCE_DIR: Final = ROOT / "docs" / "tutorial"
OUTPUT_PATH: Final = ROOT / "TabFM_Masterclass.pptx"
SOURCE_FILES: Final = (
    "01_introduction.md",
    "02_installation.md",
    "03_data_handling.md",
    "04_tabfm_mastery.md",
)

SLIDE_WIDTH: Final = 13.333
SLIDE_HEIGHT: Final = 7.5
MARGIN: Final = 0.5

BG: Final = RGBColor(15, 23, 42)
SURFACE: Final = RGBColor(30, 41, 59)
SURFACE_2: Final = RGBColor(24, 36, 55)
BORDER: Final = RGBColor(51, 65, 85)
WHITE: Final = RGBColor(255, 255, 255)
BODY: Final = RGBColor(226, 232, 240)
MUTED: Final = RGBColor(148, 163, 184)
ACCENT: Final = RGBColor(20, 184, 166)
ACCENT_DARK: Final = RGBColor(15, 118, 110)
AMBER: Final = RGBColor(245, 158, 11)
CORAL: Final = RGBColor(251, 113, 133)

FONT: Final = "Calibri"
MONO_FONT: Final = "Courier New"

EXPECTED_TITLES: Final = (
    "TabFM Streamlit Studio",
    "The Paradigm Shift",
    "Core Application Architecture",
    "Ingestion Capabilities",
    "Interactive UI Features",
    "Local Setup & Installation",
    "Performance & Pitfalls",
    "Build Your First Context",
)


@dataclass(frozen=True)
class MarkdownTable:
    """One Markdown table with normalized cell text."""

    headers: tuple[str, ...]
    rows: tuple[tuple[str, ...], ...]

    def row_for(self, label: str) -> tuple[str, ...]:
        """Return the first row whose first cell contains ``label``."""
        needle = _normalize(label)
        for row in self.rows:
            if row and needle in _normalize(row[0]):
                return row
        raise ValueError(f"Required table row not found: {label}")


@dataclass(frozen=True)
class CodeBlock:
    """A fenced Markdown code block."""

    language: str
    text: str


@dataclass
class Section:
    """Structured content belonging to one Markdown heading."""

    title: str
    level: int
    paragraphs: list[str] = field(default_factory=list)
    bullets: list[str] = field(default_factory=list)
    tables: list[MarkdownTable] = field(default_factory=list)
    code_blocks: list[CodeBlock] = field(default_factory=list)


@dataclass(frozen=True)
class MarkdownDocument:
    """Parsed tutorial document used as the presentation content source."""

    path: Path
    sections: tuple[Section, ...]
    references: dict[str, str]

    @property
    def text(self) -> str:
        parts: list[str] = []
        for section in self.sections:
            parts.extend([section.title, *section.paragraphs, *section.bullets])
            for table in section.tables:
                parts.extend(table.headers)
                parts.extend(cell for row in table.rows for cell in row)
            parts.extend(block.text for block in section.code_blocks)
        return "\n".join(parts)

    def section(self, title_fragment: str) -> Section:
        """Find a section by a stable, case-insensitive title fragment."""
        needle = _normalize(title_fragment)
        for section in self.sections:
            if needle in _normalize(section.title):
                return section
        raise ValueError(f"{self.path.name}: required section not found: {title_fragment}")

    def table(self, header_fragment: str) -> MarkdownTable:
        """Find a table whose header contains ``header_fragment``."""
        needle = _normalize(header_fragment)
        for section in self.sections:
            for table in section.tables:
                if any(needle in _normalize(header) for header in table.headers):
                    return table
        raise ValueError(f"{self.path.name}: required table not found: {header_fragment}")

    def code_containing(self, fragment: str) -> CodeBlock:
        """Find a code block containing a command fragment."""
        needle = fragment.casefold()
        for section in self.sections:
            for block in section.code_blocks:
                if needle in block.text.casefold():
                    return block
        raise ValueError(f"{self.path.name}: required code block not found: {fragment}")

    def reference_url(self, name: str) -> str:
        """Resolve a Markdown reference link by name."""
        try:
            return self.references[_normalize(name)]
        except KeyError as exc:
            raise ValueError(f"{self.path.name}: reference URL not found: {name}") from exc


def _normalize(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", value.casefold()).strip()


def _clean_inline(value: str) -> str:
    """Remove Markdown-only syntax while retaining human-readable labels."""
    value = re.sub(r"!\[([^]]*)]\([^)]+\)", r"\1", value)
    value = re.sub(r"\[([^]]+)]\([^)]+\)", r"\1", value)
    value = re.sub(r"\[([^]]+)]\[[^]]+]", r"\1", value)
    value = re.sub(r"`([^`]+)`", r"\1", value)
    value = value.replace("**", "").replace("__", "")
    value = re.sub(r"(?<!\*)\*(?!\*)", "", value)
    return re.sub(r"\s+", " ", value).strip()


def _table_cells(line: str) -> tuple[str, ...]:
    return tuple(_clean_inline(cell.strip()) for cell in line.strip().strip("|").split("|"))


def _is_table_separator(cells: Sequence[str]) -> bool:
    return bool(cells) and all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells)


def parse_markdown(path: Path) -> MarkdownDocument:
    """Parse the Markdown structures used by the four tutorial chapters."""
    if not path.is_file():
        raise FileNotFoundError(f"Required tutorial file is missing: {path}")
    raw = path.read_text(encoding="utf-8-sig")
    if not raw.strip():
        raise ValueError(f"Tutorial file is empty: {path}")

    references: dict[str, str] = {}
    sections: list[Section] = [Section(path.stem, 0)]
    current = sections[0]
    paragraph_lines: list[str] = []
    table_lines: list[str] = []
    code_lines: list[str] = []
    code_language = ""
    in_code = False

    def flush_paragraph() -> None:
        if paragraph_lines:
            text = _clean_inline(" ".join(paragraph_lines))
            if text:
                current.paragraphs.append(text)
            paragraph_lines.clear()

    def flush_table() -> None:
        if not table_lines:
            return
        parsed = [_table_cells(line) for line in table_lines]
        table_lines.clear()
        if len(parsed) < 2 or not _is_table_separator(parsed[1]):
            current.paragraphs.extend(" | ".join(row) for row in parsed if any(row))
            return
        width = len(parsed[0])
        rows = tuple(row for row in parsed[2:] if len(row) == width and any(row))
        current.tables.append(MarkdownTable(parsed[0], rows))

    for line in raw.splitlines():
        stripped = line.strip()
        if in_code:
            if stripped.startswith("```"):
                current.code_blocks.append(CodeBlock(code_language, "\n".join(code_lines).strip()))
                code_lines.clear()
                code_language = ""
                in_code = False
            else:
                code_lines.append(line.rstrip())
            continue

        reference_match = re.fullmatch(r"\[([^]]+)]:\s*(https?://\S+)", stripped)
        if reference_match:
            flush_paragraph()
            flush_table()
            references[_normalize(reference_match.group(1))] = reference_match.group(2)
            continue
        if stripped.startswith("```"):
            flush_paragraph()
            flush_table()
            in_code = True
            code_language = stripped[3:].strip().casefold()
            continue
        heading_match = re.match(r"^(#{1,6})\s+(.+)$", stripped)
        if heading_match:
            flush_paragraph()
            flush_table()
            current = Section(_clean_inline(heading_match.group(2)), len(heading_match.group(1)))
            sections.append(current)
            continue
        if stripped.startswith("|") and stripped.endswith("|"):
            flush_paragraph()
            table_lines.append(stripped)
            continue
        flush_table()
        bullet_match = re.match(r"^(?:[-*+]\s+|\d+[.)]\s+)(.+)$", stripped)
        if bullet_match:
            flush_paragraph()
            current.bullets.append(_clean_inline(bullet_match.group(1)))
            continue
        if not stripped:
            flush_paragraph()
            continue
        if stripped.startswith(">"):
            stripped = stripped.lstrip("> ")
        paragraph_lines.append(stripped)

    if in_code:
        raise ValueError(f"{path.name}: unclosed fenced code block")
    flush_paragraph()
    flush_table()
    meaningful = tuple(section for section in sections if section.level or section.paragraphs)
    if not meaningful:
        raise ValueError(f"No usable Markdown content found in {path}")
    return MarkdownDocument(path=path, sections=meaningful, references=references)


def load_sources(source_dir: Path) -> dict[str, MarkdownDocument]:
    """Load all required tutorial chapters using stable semantic keys."""
    keys = ("introduction", "installation", "data", "mastery")
    return {
        key: parse_markdown(source_dir / filename)
        for key, filename in zip(keys, SOURCE_FILES, strict=True)
    }


def _inches(value: float) -> int:
    return Inches(value)


def _estimated_lines(text: str, width: float, font_size: float) -> int:
    chars_per_line = max(8, int(width * 72 / (font_size * 0.52)))
    return sum(
        max(1, (len(line) + chars_per_line - 1) // chars_per_line) for line in text.splitlines()
    )


def _fit_font(text: str, width: float, height: float, preferred: float, minimum: float) -> float:
    size = preferred
    while size > minimum:
        capacity = max(1, int(height * 72 / (size * 1.22)))
        if _estimated_lines(text, width, size) <= capacity:
            break
        size -= 0.5
    return size


def add_text(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 14,
    minimum_size: float = 11,
    color: RGBColor = BODY,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
) -> object:
    """Add bounded text and conservatively reduce its font when necessary."""
    shape = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(width), _inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = _inches(margin)
    frame.margin_right = _inches(margin)
    frame.margin_top = _inches(margin)
    frame.margin_bottom = _inches(margin)
    fitted_size = _fit_font(text, width - 2 * margin, height - 2 * margin, size, minimum_size)
    for index, line in enumerate(text.splitlines() or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.05
        run = paragraph.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(fitted_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def add_panel(
    slide: Slide,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = SURFACE,
    line: RGBColor = BORDER,
    radius: bool = True,
) -> object:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, _inches(x), _inches(y), _inches(width), _inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_header(slide: Slide, eyebrow: str, title: str, subtitle: str | None = None) -> None:
    add_text(slide, eyebrow.upper(), MARGIN, 0.34, 5.0, 0.28, size=10, color=ACCENT, bold=True)
    add_text(
        slide, title, MARGIN, 0.66, 12.25, 0.55, size=32, minimum_size=28, color=WHITE, bold=True
    )
    if subtitle:
        add_text(slide, subtitle, MARGIN, 1.25, 12.1, 0.42, size=14, color=MUTED)


def add_footer(slide: Slide, number: int, source: str) -> None:
    add_text(slide, source, MARGIN, 7.13, 10.8, 0.18, size=8.5, color=MUTED)
    add_text(
        slide,
        f"{number:02d}",
        12.27,
        7.08,
        0.55,
        0.22,
        size=9,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_badge(
    slide: Slide,
    label: str,
    x: float,
    y: float,
    width: float,
    *,
    fill: RGBColor = ACCENT_DARK,
    color: RGBColor = WHITE,
) -> None:
    add_panel(slide, x, y, width, 0.35, fill=fill, line=fill)
    add_text(
        slide,
        label,
        x + 0.08,
        y + 0.04,
        width - 0.16,
        0.24,
        size=10,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_arrow(slide: Slide, x1: float, y: float, x2: float) -> None:
    arrow = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, _inches(x1), _inches(y), _inches(x2), _inches(y)
    )
    arrow.line.color.rgb = ACCENT
    arrow.line.width = Pt(2.5)
    chevron = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        _inches(x2 - 0.13),
        _inches(y - 0.12),
        _inches(0.2),
        _inches(0.24),
    )
    chevron.fill.solid()
    chevron.fill.fore_color.rgb = ACCENT
    chevron.line.fill.background()


def _new_slide(prs: PresentationType) -> Slide:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = BG
    return slide


def _comparison_value(table: MarkdownTable, label: str, column: int) -> str:
    row = table.row_for(label)
    if len(row) <= column:
        raise ValueError(f"Comparison row '{label}' is missing column {column + 1}")
    return row[column]


def build_title_slide(prs: PresentationType, intro: MarkdownDocument) -> None:
    slide = _new_slide(prs)
    intro.section("What is a tabular foundation model")
    tagline = "A local-first path from labeled context rows to zero-shot tabular predictions."

    add_badge(slide, "ZERO TO MASTER COURSE", 0.58, 0.62, 2.25)
    add_text(
        slide,
        "TabFM\nStreamlit Studio",
        0.58,
        1.25,
        6.1,
        1.7,
        size=38,
        minimum_size=34,
        color=WHITE,
        bold=True,
    )
    add_text(slide, tagline, 0.62, 3.18, 5.55, 0.92, size=18, minimum_size=16, color=BODY)
    add_text(
        slide,
        "Google Research TabFM · Classification + Regression",
        0.62,
        4.34,
        5.8,
        0.38,
        size=12,
        color=ACCENT,
        bold=True,
    )

    add_panel(slide, 7.03, 0.72, 5.65, 5.95, fill=SURFACE_2, line=RGBColor(51, 65, 85))
    add_text(
        slide,
        "ONE TABLE. ONE CONTEXT. ONE PREDICTION.",
        7.42,
        1.08,
        4.9,
        0.35,
        size=11,
        color=MUTED,
        bold=True,
    )
    headers = ("age", "plan", "visits", "target")
    rows = (("34", "pro", "12", "retain"), ("51", "basic", "3", "churn"), ("42", "pro", "8", "?"))
    cell_w, cell_h = 1.13, 0.7
    origin_x, origin_y = 7.45, 1.75
    for col, header in enumerate(headers):
        add_panel(
            slide, origin_x + col * cell_w, origin_y, 1.03, 0.52, fill=ACCENT_DARK, line=ACCENT_DARK
        )
        add_text(
            slide,
            header,
            origin_x + col * cell_w,
            origin_y + 0.13,
            1.03,
            0.22,
            size=10,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    for row_index, row in enumerate(rows):
        for col, value in enumerate(row):
            highlight = row_index == 2 and col == 3
            fill = ACCENT if highlight else SURFACE
            text_color = BG if highlight else BODY
            add_panel(
                slide,
                origin_x + col * cell_w,
                origin_y + 0.65 + row_index * cell_h,
                1.03,
                0.58,
                fill=fill,
                line=RGBColor(71, 85, 105),
            )
            add_text(
                slide,
                value,
                origin_x + col * cell_w,
                origin_y + 0.82 + row_index * cell_h,
                1.03,
                0.22,
                size=11,
                color=text_color,
                bold=highlight,
                align=PP_ALIGN.CENTER,
            )
    add_arrow(slide, 8.12, 4.75, 11.78)
    add_text(
        slide,
        "Frozen model conditions on labeled rows",
        7.55,
        4.92,
        4.9,
        0.34,
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_panel(slide, 8.42, 5.48, 2.9, 0.72, fill=ACCENT, line=ACCENT)
    add_text(
        slide,
        "PREDICTION  retain",
        8.58,
        5.68,
        2.58,
        0.26,
        size=13,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 1, "Source: 01_introduction.md")


def build_paradigm_slide(prs: PresentationType, intro: MarkdownDocument) -> None:
    slide = _new_slide(prs)
    table = intro.table("Task-specific GBDT")
    add_header(
        slide,
        "FOUNDATIONS",
        "The Paradigm Shift",
        "Task-specific optimization versus frozen in-context adaptation",
    )

    columns = (
        (0.62, "TRADITIONAL GBDT", "XGBoost / LightGBM", CORAL, 1),
        (6.98, "PRETRAINED TABFM", "Google Research TabFM", ACCENT, 2),
    )
    labels = ("Adaptation mechanism", "Per-dataset weight updates", "Main cost")
    for x, eyebrow, title, accent, column in columns:
        add_panel(slide, x, 1.88, 5.72, 4.72, fill=SURFACE_2, line=accent)
        add_text(slide, eyebrow, x + 0.35, 2.16, 4.9, 0.26, size=10, color=accent, bold=True)
        add_text(slide, title, x + 0.35, 2.48, 4.9, 0.42, size=22, color=WHITE, bold=True)
        for index, label in enumerate(labels):
            y = 3.18 + index * 1.02
            add_text(slide, label.upper(), x + 0.35, y, 2.2, 0.22, size=9, color=MUTED, bold=True)
            value = _comparison_value(table, label, column)
            add_text(
                slide, value, x + 0.35, y + 0.27, 4.95, 0.58, size=13, minimum_size=11.5, color=BODY
            )
        footer = "Mature CPU tooling" if column == 1 else "Immediate zero-shot baseline"
        add_badge(
            slide, footer, x + 0.35, 5.95, 2.35, fill=accent, color=BG if column == 2 else WHITE
        )
    add_panel(slide, 6.18, 3.48, 0.96, 0.96, fill=BG, line=ACCENT)
    add_text(
        slide, "VS", 6.28, 3.76, 0.76, 0.25, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER
    )
    add_footer(slide, 2, "Source: 01_introduction.md · comparison table")


def build_architecture_slide(prs: PresentationType, intro: MarkdownDocument) -> None:
    slide = _new_slide(prs)
    intro.section("Preprocessing before the neural model")
    intro.section("Classification and regression")
    add_header(
        slide,
        "SYSTEM FLOW",
        "Core Application Architecture",
        "No task-specific weight updates: labeled examples become inference context",
    )

    stages = (
        ("01", "DATA INGESTION", "Files, HTTPS, Kaggle, Hugging Face"),
        ("02", "CONTEXT PREPARATION", "Validate schema, encode mixed types, select labeled rows"),
        ("03", "FROZEN INFERENCE", "Eight ensemble views condition on context + test rows"),
        ("04", "PREDICTOR UI", "Labels, probabilities, values, metrics, CSV export"),
    )
    x_positions = (0.58, 3.78, 6.98, 10.18)
    for index in range(len(stages) - 1):
        add_arrow(slide, x_positions[index] + 2.55, 3.43, x_positions[index + 1] - 0.13)
    for x, (number, title, description) in zip(x_positions, stages, strict=True):
        add_panel(slide, x, 2.16, 2.55, 2.55, fill=SURFACE_2, line=RGBColor(51, 65, 85))
        add_badge(slide, number, x + 0.2, 2.38, 0.55)
        add_text(
            slide,
            title,
            x + 0.2,
            2.92,
            2.15,
            0.58,
            size=14,
            minimum_size=12,
            color=WHITE,
            bold=True,
        )
        add_text(
            slide, description, x + 0.2, 3.62, 2.15, 0.76, size=11.5, minimum_size=10.5, color=BODY
        )

    add_panel(slide, 2.27, 5.18, 4.06, 1.12, fill=SURFACE, line=ACCENT_DARK)
    add_text(slide, "CLASSIFICATION", 2.56, 5.44, 1.7, 0.24, size=11, color=ACCENT, bold=True)
    add_text(
        slide, "2–10 classes · labels + probabilities", 2.56, 5.76, 3.35, 0.27, size=12, color=BODY
    )
    add_panel(slide, 7.0, 5.18, 4.06, 1.12, fill=SURFACE, line=ACCENT_DARK)
    add_text(slide, "REGRESSION", 7.29, 5.44, 1.55, 0.24, size=11, color=ACCENT, bold=True)
    add_text(
        slide, "Numeric target · continuous prediction", 7.29, 5.76, 3.35, 0.27, size=12, color=BODY
    )
    add_footer(slide, 3, "Sources: 01_introduction.md · 04_tabfm_mastery.md")


def _require_terms(document: MarkdownDocument, terms: Iterable[str]) -> None:
    normalized = _normalize(document.text)
    missing = [term for term in terms if _normalize(term) not in normalized]
    if missing:
        raise ValueError(
            f"{document.path.name}: required content terms missing: {', '.join(missing)}"
        )


def build_ingestion_slide(prs: PresentationType, data_doc: MarkdownDocument) -> None:
    slide = _new_slide(prs)
    _require_terms(data_doc, ("CSV", "Parquet", "XLSX", "Kaggle", "Hugging Face", "HTTPS"))
    add_header(
        slide,
        "DATA ENGINEERING",
        "Ingestion Capabilities",
        "Multiple sources converge on one validated in-memory DataFrame",
    )

    groups = (
        (
            0.62,
            "LOCAL FILES",
            (
                ("CSV", "Delimited tables"),
                ("PARQUET", "Typed columnar data"),
                ("XLSX", "Excel workbooks"),
            ),
        ),
        (
            6.72,
            "CONNECTED SOURCES",
            (
                ("HTTPS", "Bounded direct download"),
                ("KAGGLE", "Search + dataset pull"),
                ("HF HUB", "Repository file download"),
            ),
        ),
    )
    for x, title, cards in groups:
        add_panel(slide, x, 1.9, 5.98, 3.92, fill=SURFACE_2, line=RGBColor(51, 65, 85))
        add_text(slide, title, x + 0.34, 2.18, 4.9, 0.28, size=11, color=ACCENT, bold=True)
        for index, (label, detail) in enumerate(cards):
            card_x = x + 0.34 + index * 1.83
            add_panel(slide, card_x, 2.75, 1.6, 2.42, fill=SURFACE, line=ACCENT_DARK)
            add_panel(slide, card_x + 0.2, 3.04, 0.54, 0.54, fill=ACCENT_DARK, line=ACCENT_DARK)
            add_text(
                slide,
                str(index + 1),
                card_x + 0.2,
                3.18,
                0.54,
                0.2,
                size=11,
                color=WHITE,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
            add_text(slide, label, card_x + 0.2, 3.82, 1.2, 0.28, size=14, color=WHITE, bold=True)
            add_text(
                slide, detail, card_x + 0.2, 4.24, 1.2, 0.58, size=11, minimum_size=9.5, color=BODY
            )

    add_panel(slide, 0.62, 6.06, 12.08, 0.65, fill=RGBColor(17, 48, 55), line=ACCENT_DARK)
    add_text(slide, "LOCAL-FIRST SECURITY", 0.9, 6.27, 2.0, 0.22, size=10, color=ACCENT, bold=True)
    add_text(
        slide,
        "Environment-only credentials · HTTPS by default · bounded downloads · "
        "dedicated workspace paths",
        3.0,
        6.22,
        9.15,
        0.3,
        size=11.5,
        color=BODY,
    )
    add_footer(slide, 4, "Source: 03_data_handling.md")


def build_ui_slide(prs: PresentationType, mastery: MarkdownDocument) -> None:
    slide = _new_slide(prs)
    widgets = mastery.table("Context dtype/shape")
    result_table = mastery.table("Element")
    widgets.row_for("Numeric")
    widgets.row_for("Categorical")
    result_table.row_for("Inference latency")
    result_table.row_for("Probability columns")
    add_header(
        slide,
        "STREAMLIT EXPERIENCE",
        "Interactive UI Features",
        "Typed manual inputs and immediate, inspectable prediction results",
    )

    add_panel(slide, 0.62, 1.85, 5.17, 4.96, fill=SURFACE_2, line=RGBColor(51, 65, 85))
    add_text(slide, "SINGLE TEST CASE", 0.97, 2.14, 3.2, 0.25, size=11, color=ACCENT, bold=True)
    form_rows = (("age", "42.0"), ("plan", "pro    v"), ("active", "True   v"))
    for index, (label, value) in enumerate(form_rows):
        y = 2.72 + index * 0.86
        add_text(slide, label, 0.98, y, 1.0, 0.23, size=10, color=MUTED, bold=True)
        add_panel(slide, 1.95, y - 0.06, 3.36, 0.54, fill=BG, line=RGBColor(71, 85, 105))
        add_text(slide, value, 2.12, y + 0.08, 2.98, 0.22, size=11.5, color=BODY)
    add_panel(slide, 1.0, 5.43, 4.29, 0.62, fill=ACCENT, line=ACCENT)
    add_text(
        slide,
        "PREDICT SINGLE CASE",
        1.12,
        5.62,
        4.05,
        0.23,
        size=12,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Widgets follow the prepared context schema",
        1.0,
        6.27,
        4.3,
        0.25,
        size=10.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_panel(slide, 6.08, 1.85, 6.62, 4.96, fill=SURFACE_2, line=ACCENT_DARK)
    add_text(
        slide, "ILLUSTRATIVE RESULT PANEL", 6.42, 2.14, 3.2, 0.25, size=11, color=ACCENT, bold=True
    )
    add_panel(slide, 6.42, 2.62, 2.27, 1.14, fill=ACCENT, line=ACCENT)
    add_text(slide, "PREDICTION", 6.68, 2.83, 1.75, 0.21, size=9, color=BG, bold=True)
    add_text(slide, "retain", 6.68, 3.13, 1.75, 0.32, size=20, color=BG, bold=True)
    metric_cards = (("ACCURACY", "0.91"), ("LATENCY", "84 ms"))
    for index, (label, value) in enumerate(metric_cards):
        x = 8.93 + index * 1.62
        add_panel(slide, x, 2.62, 1.42, 1.14, fill=SURFACE, line=RGBColor(71, 85, 105))
        add_text(slide, label, x + 0.15, 2.84, 1.1, 0.18, size=8.5, color=MUTED, bold=True)
        add_text(slide, value, x + 0.15, 3.14, 1.1, 0.27, size=16, color=WHITE, bold=True)
    probabilities = (("retain", 0.78, ACCENT), ("churn", 0.22, CORAL))
    for index, (label, probability, color) in enumerate(probabilities):
        y = 4.13 + index * 0.67
        add_text(slide, label, 6.42, y, 1.0, 0.23, size=10.5, color=BODY)
        add_panel(slide, 7.46, y + 0.03, 3.8, 0.2, fill=BG, line=BG, radius=False)
        add_panel(
            slide, 7.46, y + 0.03, 3.8 * probability, 0.2, fill=color, line=color, radius=False
        )
        add_text(
            slide,
            f"{probability:.0%}",
            11.48,
            y - 0.01,
            0.64,
            0.23,
            size=10.5,
            color=BODY,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
    add_panel(slide, 6.42, 5.62, 5.86, 0.62, fill=SURFACE, line=RGBColor(71, 85, 105))
    add_text(
        slide,
        "Warnings  ·  Data table  ·  Download predictions.csv",
        6.7,
        5.82,
        5.3,
        0.23,
        size=11,
        color=BODY,
    )
    add_text(
        slide,
        "Batch mode uses the same aligned result contract",
        6.42,
        6.36,
        5.8,
        0.24,
        size=10.5,
        color=MUTED,
    )
    add_footer(slide, 5, "Source: 04_tabfm_mastery.md · manual case and result panel")


def _command_lines(block: CodeBlock, *, include: Sequence[str] | None = None) -> list[str]:
    lines = [line.rstrip() for line in block.text.splitlines() if line.strip()]
    if include is None:
        return lines
    return [line for line in lines if any(token.casefold() in line.casefold() for token in include)]


def build_setup_slide(prs: PresentationType, installation: MarkdownDocument) -> None:
    slide = _new_slide(prs)
    clone = installation.code_containing("git clone")
    sync = installation.code_containing("uv sync --locked")
    env = installation.code_containing("Copy-Item .env.example .env")
    run = installation.code_containing("streamlit run app.py")
    add_header(
        slide,
        "DEVELOPER QUICKSTART",
        "Local Setup & Installation",
        "Three PowerShell steps with a reproducible uv environment",
    )

    commands = (
        (
            "01",
            "CLONE",
            _command_lines(clone, include=("git clone", "Set-Location google")),
        ),
        (
            "02",
            "ENVIRONMENT",
            ["uv python install 3.12.10", *_command_lines(sync, include=("uv sync",))[-1:]],
        ),
        (
            "03",
            "CONFIGURE + RUN",
            [*_command_lines(env)[:1], "# Review license, then set .env", *_command_lines(run)[:1]],
        ),
    )
    for index, (number, title, lines) in enumerate(commands):
        x = 0.62 + index * 4.08
        add_panel(slide, x, 1.93, 3.72, 4.72, fill=SURFACE_2, line=RGBColor(51, 65, 85))
        add_badge(slide, number, x + 0.3, 2.24, 0.58)
        add_text(slide, title, x + 1.02, 2.29, 2.3, 0.3, size=14, color=WHITE, bold=True)
        add_panel(
            slide, x + 0.3, 3.02, 3.12, 2.72, fill=RGBColor(8, 15, 29), line=RGBColor(51, 65, 85)
        )
        add_text(
            slide,
            "PS >",
            x + 0.5,
            3.27,
            0.48,
            0.22,
            size=10,
            color=ACCENT,
            bold=True,
            font=MONO_FONT,
        )
        add_text(
            slide,
            "\n\n".join(lines),
            x + 0.5,
            3.72,
            2.72,
            1.58,
            size=11.5,
            minimum_size=9.5,
            color=BODY,
            font=MONO_FONT,
        )
        caption = (
            "Repository + tutorial",
            "CPU-safe default; choose CUDA when available",
            "Local secrets stay outside Git",
        )
        add_text(
            slide,
            caption[index],
            x + 0.34,
            5.98,
            3.02,
            0.33,
            size=10.5,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
    add_footer(slide, 6, "Source: 02_installation.md · PowerShell path")


def build_performance_slide(
    prs: PresentationType, intro: MarkdownDocument, mastery: MarkdownDocument
) -> None:
    slide = _new_slide(prs)
    _require_terms(intro, ("2–10 classes", "500 features", "6.56 GB"))
    _require_terms(mastery, ("5,000", "Fixed at 8", "CPU appears hung", "Missing-column"))
    add_header(
        slide,
        "OPERATING ENVELOPE",
        "Performance & Pitfalls",
        "Scale context deliberately and evaluate on representative held-out data",
    )

    metrics = (
        ("2–10", "CLASSIFICATION\nCLASSES"),
        ("≤500", "FEATURE\nCOLUMNS"),
        ("≤5,000", "CONTEXT ROWS\nPER MEMBER"),
        ("8", "ENSEMBLE\nMEMBERS"),
    )
    for index, (value, label) in enumerate(metrics):
        x = 0.62 + index * 3.03
        add_panel(slide, x, 1.86, 2.72, 1.68, fill=SURFACE_2, line=ACCENT_DARK)
        add_text(
            slide,
            value,
            x + 0.25,
            2.15,
            2.22,
            0.48,
            size=26,
            color=ACCENT,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            label,
            x + 0.25,
            2.77,
            2.22,
            0.48,
            size=10,
            color=BODY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    warnings = (
        (AMBER, "CHECKPOINT WEIGHT", "~6.6 GB per task plus cache headroom"),
        (CORAL, "CPU LATENCY", "Functional fallback; transformer inference may be very slow"),
        (AMBER, "CONTEXT STATE", "Prepare labeled context before batch or single prediction"),
        (
            CORAL,
            "SCHEMA DRIFT",
            "Missing columns are filled; extra columns are ignored with warnings",
        ),
        (AMBER, "LICENSE BOUNDARY", "Non-commercial research and evaluation only"),
        (CORAL, "MEMORY PRESSURE", "Reduce context rows, features, or test batch size first"),
    )
    for index, (color, title, detail) in enumerate(warnings):
        col, row = index % 2, index // 2
        x, y = 0.62 + col * 6.08, 3.91 + row * 0.89
        add_panel(slide, x, y, 5.74, 0.7, fill=SURFACE, line=RGBColor(51, 65, 85))
        add_panel(slide, x + 0.18, y + 0.18, 0.28, 0.28, fill=color, line=color)
        add_text(slide, title, x + 0.62, y + 0.12, 1.72, 0.2, size=9, color=color, bold=True)
        add_text(
            slide, detail, x + 0.62, y + 0.34, 4.75, 0.24, size=10.5, minimum_size=9.5, color=BODY
        )
    add_footer(slide, 7, "Sources: 01_introduction.md · 04_tabfm_mastery.md")


def _project_url(installation: MarkdownDocument) -> str:
    clone = installation.code_containing("git clone").text
    match = re.search(r"https://github\.com/[\w.-]+/[\w.-]+(?:\.git)?", clone)
    if not match:
        raise ValueError("Could not extract the project GitHub URL from installation instructions")
    return match.group(0).removesuffix(".git")


def add_link_card(
    slide: Slide, title: str, label: str, url: str, x: float, y: float, width: float
) -> None:
    panel = add_panel(slide, x, y, width, 1.28, fill=SURFACE, line=ACCENT_DARK)
    panel.click_action.hyperlink.address = url
    add_text(slide, title, x + 0.24, y + 0.18, width - 0.48, 0.25, size=10, color=ACCENT, bold=True)
    add_text(
        slide,
        label,
        x + 0.24,
        y + 0.53,
        width - 0.48,
        0.5,
        size=10.5,
        minimum_size=9,
        color=WHITE,
        bold=True,
    )


def build_cta_slide(
    prs: PresentationType, intro: MarkdownDocument, installation: MarkdownDocument
) -> None:
    slide = _new_slide(prs)
    project_url = _project_url(installation)
    docs_url = f"{project_url}#zero-to-master-tutorial"
    docs_label = docs_url.removeprefix("https://")
    official_url = intro.reference_url("tabfm-repo")

    add_badge(slide, "NEXT STEP", 0.62, 0.62, 1.2)
    add_text(
        slide, "Build Your First Context", 0.62, 1.22, 8.2, 0.68, size=34, color=WHITE, bold=True
    )
    add_text(
        slide,
        "Clone the workbench, accept the model terms, and turn a labeled table into an "
        "inspectable prediction.",
        0.62,
        2.02,
        7.85,
        0.78,
        size=17,
        minimum_size=15,
        color=BODY,
    )

    steps = (("01", "CLONE"), ("02", "PREPARE CONTEXT"), ("03", "PREDICT + EVALUATE"))
    for index, (number, label) in enumerate(steps):
        x = 0.62 + index * 2.75
        add_panel(slide, x, 3.18, 2.45, 0.78, fill=SURFACE_2, line=RGBColor(51, 65, 85))
        add_text(slide, number, x + 0.18, 3.39, 0.42, 0.22, size=10, color=ACCENT, bold=True)
        add_text(slide, label, x + 0.66, 3.38, 1.6, 0.23, size=10.5, color=WHITE, bold=True)

    add_panel(slide, 9.0, 0.72, 3.68, 5.18, fill=SURFACE_2, line=ACCENT_DARK)
    add_text(slide, "RESOURCES", 9.35, 1.07, 2.4, 0.26, size=11, color=ACCENT, bold=True)
    add_link_card(
        slide,
        "PROJECT",
        "github.com/pypi-ahmad/google-tabfm-project",
        project_url,
        9.35,
        1.62,
        2.98,
    )
    add_link_card(
        slide,
        "DOCUMENTATION",
        docs_label,
        docs_url,
        9.35,
        3.04,
        2.98,
    )
    add_link_card(
        slide,
        "UPSTREAM",
        "github.com/google-research/tabfm",
        official_url,
        9.35,
        4.46,
        2.98,
    )
    add_panel(slide, 0.62, 4.55, 7.78, 1.35, fill=RGBColor(17, 48, 55), line=ACCENT_DARK)
    add_text(
        slide,
        "RESEARCH-ONLY MODEL WEIGHTS",
        0.95,
        4.85,
        3.0,
        0.26,
        size=11,
        color=ACCENT,
        bold=True,
    )
    add_text(
        slide,
        "Review and acknowledge the TabFM Non-Commercial License before loading a checkpoint.",
        0.95,
        5.24,
        6.92,
        0.34,
        size=13,
        color=BODY,
    )
    add_footer(slide, 8, "Sources: 01_introduction.md · 02_installation.md")


def _all_text(slide: Slide) -> str:
    return "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))


def validate_presentation(path: Path) -> None:
    """Reopen the deck and validate its structural and content contract."""
    prs = Presentation(path)
    if len(prs.slides) != 8:
        raise RuntimeError(f"Expected 8 slides, found {len(prs.slides)}")
    if prs.slide_width != _inches(SLIDE_WIDTH) or prs.slide_height != _inches(SLIDE_HEIGHT):
        raise RuntimeError("Presentation is not 16:9 widescreen")

    hyperlink_count = 0
    for index, (slide, title) in enumerate(zip(prs.slides, EXPECTED_TITLES, strict=True), start=1):
        slide_text = _all_text(slide)
        if _normalize(title) not in _normalize(slide_text):
            raise RuntimeError(f"Slide {index} is missing its required title: {title}")
        if "TODO" in slide_text.upper():
            raise RuntimeError(f"Slide {index} contains placeholder text")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                raise RuntimeError(f"Slide {index} has a shape outside the top/left boundary")
            if shape.left + shape.width > prs.slide_width + 1:
                raise RuntimeError(f"Slide {index} has a shape outside the right boundary")
            if shape.top + shape.height > prs.slide_height + 1:
                raise RuntimeError(f"Slide {index} has a shape outside the bottom boundary")
            if shape.click_action.hyperlink.address:
                hyperlink_count += 1
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    hyperlink_count += sum(1 for run in paragraph.runs if run.hyperlink.address)
    if hyperlink_count < 3:
        raise RuntimeError(f"Expected at least 3 resource hyperlinks, found {hyperlink_count}")


def generate_presentation(source_dir: Path, output_path: Path) -> Path:
    """Compile tutorial content into an atomic, validated PowerPoint file."""
    documents = load_sources(source_dir)
    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = _inches(SLIDE_WIDTH)
    prs.slide_height = _inches(SLIDE_HEIGHT)
    build_title_slide(prs, documents["introduction"])
    build_paradigm_slide(prs, documents["introduction"])
    build_architecture_slide(prs, documents["introduction"])
    build_ingestion_slide(prs, documents["data"])
    build_ui_slide(prs, documents["mastery"])
    build_setup_slide(prs, documents["installation"])
    build_performance_slide(prs, documents["introduction"], documents["mastery"])
    build_cta_slide(prs, documents["introduction"], documents["installation"])

    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix=f".{output_path.stem}-", suffix=".pptx", dir=output_path.parent, delete=False
        ) as temporary:
            temporary_path = Path(temporary.name)
        prs.save(temporary_path)
        validate_presentation(temporary_path)
        temporary_path.replace(output_path)
    finally:
        if temporary_path is not None and temporary_path.exists():
            temporary_path.unlink()
    return output_path


def main() -> int:
    """Generate the deck using repository-relative defaults."""
    try:
        generated = generate_presentation(SOURCE_DIR, OUTPUT_PATH)
    except (OSError, ValueError, RuntimeError) as exc:
        print(f"Failed to create presentation: {exc}", file=sys.stderr)
        return 1
    print(f"Created {generated.name} successfully.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
